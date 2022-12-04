'''
Módulo para leitura de apresentações PowerPoint.
'''
import collections.abc
from pptx import Presentation
from pathlib import Path
from module.comparator import data_finder

def ppt_r(file_path, word_list):
    '''Função para abertura e leitura de apresentações PowerPoint'''
    temp, laux = dict(), list()
    prs = Presentation(file_path)
    for n_slide, slide in enumerate(prs.slides, start = 1):
        for shape in slide.shapes:
            if shape.has_text_frame and len(shape.text) > 0:
                laux = data_finder(word_list, shape.text, n_slide, file_path.suffix)
                if laux:
                    temp[n_slide].append(laux[n_slide][0]) if n_slide in temp else temp.update(laux)
    if len(temp) > 0:
        return temp

if __name__ == '__main__':
    pass

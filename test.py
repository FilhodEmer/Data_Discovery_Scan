import pptx
from glob import glob

for file in glob('*.PPTX'):
    prs = pptx.Presentation(file)
    aux = dict()
    for i, slide in enumerate(prs.slides, start = 1):
        for shape in slide.shapes:
            boolean = False
            if shape.has_text_frame and 'dado' in shape.text:
                boolean, key = True, i
                if key in aux:
                    aux[key] += ';{}'.format(shape.text)
                    continue
                aux[key] = shape.text
    print('Arquivo: {}'.format(file)) if boolean == True else ()
    for x, y in aux.items():
        print('Slide {}\n{}'.format(x, y.split('\n'))) if 'dado' in y else ()


            #print('Arquivo: {}\nSlide {}\n{}'.format(file, i, shape.text)) if shape.has_text_frame and 'dado' in shape.text else ()

'''for fname in glob ('*.pptx'):
    print('Arquivo: ', fname, '\n')
    prs = pptx.Presentation(fname)

    for i, sld in enumerate(prs.slides, start=1):

        print(f'-- Slide {i} --')

        for shp in sld.shapes:
            
            if shp.has_text_frame:
                print(shp.text.strip())

            if shp.has_table:
                tbl = shp.table
                row_count, col_count = len(tbl.rows), len(tbl.columns)
                #row_count = len(tbl.rows)
                #col_count = len(tbl.columns)
                for r in range(0, row_count):                 
                    text=''
                    for c in range(0, col_count):
                        cell = tbl.cell(r, c)
                        paragraphs = cell.text_frame.paragraphs 
                        for paragraph in paragraphs:
                            for run in paragraph.runs:
                                text += run.text
                            text += ', '
                    print(text)
            print()'''
from re import I
from docx2txt import process
from pathlib import Path
from time import time
from pptx import Presentation
from xlrd import open_workbook_xls
from openpyxl import load_workbook
import fitz

start = time()
list2 = list(map(chr, range(97, 123)))
with open('f:/Desktop/Saida_Preliminar.txt', 'w', encoding = 'UTF-8') as out:
    #for disk in list2:
    for file in list(Path(r'c:\Users\emers\Desktop\Faculdade\TCC\Scan_DataDiscovery\teste').rglob('*.*')):
        #out.write(str(file))
        #out.write('\n')
        aux, counter = dict(), int()
        if '~$' not in file.name and not Path.stat(file).st_size == 0:
            try:
                if file.suffix in ['.docx', '.DOCX']:
                    document = process(file)
                    document.strip()
                    list_document = document.split(sep = '\n')
                    while '' in list_document:
                        temp = list_document.index('')
                        del(list_document[temp])
                    for line in range(len(list_document)):
                        key = line
                        if 'dado' in list_document[line]:
                            counter += 1
                            aux[key] = list_document[line].strip()
                    if counter > 0:
                        out.write('Arquivo: {}\n'.format(file))
                        for num, cont in aux.items():
                            out.write('Linha {}: {}\n'.format(num, cont))
                        out.write('\n')
                
                elif file.suffix in ['.pdf', '.PDF']:
                    for page in fitz.open(file):
                        content = str()
                        content += page.get_text()
                        counter += 1
                        if 'dado' in content:
                            key = counter
                            aux[key] = content
                    if len(aux) > 0:
                        out.write('Arquivo: {}\n'.format(file))
                        for num in aux.keys():
                            out.write('Página {}\n'.format(num))
                        out.write('\n')
                
                elif file.suffix in ['.pptx', '.PPTX']:
                    laux = list()
                    prs = Presentation(file)
                    for i, slide in enumerate(prs.slides, start = 1):
                        for shape in slide.shapes:
                            if shape.has_text_frame and 'dado' in shape.text:
                                counter += 1
                                key = file
                                laux.append([i, shape.text]) if len(laux) == 0 or i not in laux[-1] else laux[-1].append(shape.text)
                    if counter > 0:
                        aux[key] = laux
                        out.write('Arquivo: {}\n'.format(key))
                        for cont in aux.values():
                            for i in range(len(cont)):
                                out.write('Slide ')
                                for x in range(len(cont[i])):
                                    out.write('{}\n'.format(cont[i][x]))
                
                elif file.suffix in ['.xlsx', '.xls', '.xlsm', '.XLSX', '.XLS', '.XLSM']:
                    if file.suffix in ['.xls', '.XLS']:
                        df = open_workbook_xls(file)
                        for x in range(df.nsheets):
                            sheet = df.sheet_by_index(x)
                            for i in range(sheet.nrows):
                                for j in range(sheet.ncols):
                                    if 'dado' in str(sheet.cell_value(i, j)):
                                        counter += 1
                                        key = sheet.name
                                        aux[key] = str(chr(j + 97).upper()) + str(i + 1), sheet.cell_value(i, j)
                            if counter > 0:
                                out.write('Arquivo: {}\n'.format(file))
                                out.write('Aba: {}\n'.format(key))
                                for cont in aux.values():
                                    out.write('Célula: {}\n{}\n'.format(cont[0], cont[1], end = '\t'))
                    
                    elif file.suffix in ['.xlsx', '.XLSX']:
                        dataframe = load_workbook(file)
                        for sheet in dataframe.worksheets:
                            sheet.active_cell
                            for row in range(0, sheet.max_row):
                                hlp = int()
                                for col in sheet.iter_cols(1, sheet.max_column):
                                    hlp += 1
                                    cell = str(col[row].value)
                                    if 'dado' in cell:
                                        counter += 1
                                        key = sheet.title
                                        aux[key] = str(chr(hlp + 96).upper()) + str(row + 1), cell
                            if counter > 0:
                                out.write('Arquivo: {}\n'.format(file))
                                out.write('Aba: {}\n'.format(key))
                                for cont in aux.values():
                                    out.write('Célula {}: {}\n'.format(cont[0], cont[1]))
                
                else:
                    if file.name != 'Saida_Preliminar.txt':
                        count = int()
                        for line in open(r'file', encoding = 'UTF-8'):
                            count += 1
                            if 'dado' in line.strip():
                                key = count
                                counter += 1
                                aux[key] = line.strip()
                    if counter > 0:
                        out.write('Arquivo: {}\n'.format(file))
                        for num, cont in aux.items():
                            out.write('Linha {}: {}\n'.format(num, cont))
                        out.write('\n')
            except (UnicodeDecodeError, PermissionError, FileNotFoundError, ValueError):
                continue

print('Saída Gravada.')
end = time()

print('{:0.3f}s'.format(end - start))